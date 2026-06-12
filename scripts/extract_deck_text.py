#!/usr/bin/env python3
"""Dump a deck's on-slide text (and optionally its speaker notes), walking groups.

Used in two places in the method:
  * REPAIR track, citations step -- read the slide text to compare each claim against
    the source and find dropped attributions.
  * BUILD Phase 3 QA / REPAIR notes & rebuild steps -- read the actual slide
    order/content so the spoken script and the rebuild map to what's really on
    each slide.

This replaces the copy-paste extraction snippet that used to live inline in SKILL.md.

Usage:
  python3 extract_deck_text.py --deck deck.pptx                # text only, human readable
  python3 extract_deck_text.py --deck deck.pptx --notes        # also dump speaker notes
  python3 extract_deck_text.py --deck deck.pptx --json         # {"1": {...}} for tooling
"""
import argparse
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def walk(shapes):
    for s in shapes:
        if s.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk(s.shapes)
        else:
            yield s


def slide_lines(slide):
    out = []
    for sh in walk(slide.shapes):
        if sh.has_text_frame and sh.text_frame.text.strip():
            out.append(sh.text_frame.text.replace("\n", " | "))
    return out


def main():
    ap = argparse.ArgumentParser(description="Extract slide text/notes from a .pptx.")
    ap.add_argument("--deck", required=True)
    ap.add_argument("--notes", action="store_true", help="include speaker notes")
    ap.add_argument("--json", action="store_true", help="emit JSON instead of text")
    args = ap.parse_args()

    p = Presentation(args.deck)
    data = {}
    for i, sl in enumerate(p.slides, 1):
        lines = slide_lines(sl)
        notes = ""
        if args.notes and sl.has_notes_slide:
            notes = sl.notes_slide.notes_text_frame.text.strip()
        data[str(i)] = {"text": lines, "notes": notes}

    if args.json:
        print(json.dumps(data, ensure_ascii=False, indent=2))
        return

    for k, v in data.items():
        print(f"==== SLIDE {k} ====")
        for line in v["text"]:
            print(line)
        if args.notes:
            print(f"  -- notes ({len(v['notes'].split())} words) --")
            print(f"  {v['notes']}")
        print()


if __name__ == "__main__":
    main()

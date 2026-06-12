#!/usr/bin/env python3
"""Reinstate dropped inline citations in a .pptx (Repair track).

Locates each target paragraph by a UNIQUE substring within a given slide, then appends
the citation to the paragraph's last run -- before the final period if present -- so the
new text inherits the run's formatting.

Edits are authored in a JSON file (consistent with the "author text in JSON" rule), so
you never hand-edit this script. Each edit is [slide_number, unique_substring, citation]:

    [
      [2, "misaligned, generic, and inauthentic content", "(Sun et al., 2025)"],
      [6, "Prevents generic, off-topic superficiality",    "(Kulkarni, 2024)"]
    ]

Safety:
  * Dry run by default; --apply writes.
  * Idempotent: a paragraph that already contains the citation is skipped (re-running
    the same edits never double-appends).
  * Auto-backup before apply to <deck>_BACKUP_citations.pptx; an existing backup is
    never overwritten (a timestamped name is used instead).

Workflow:
  1. Author edits in citations.json (see references/citations_template.json).
  2. Dry run:   python3 reinstate_citations.py --deck deck.pptx --edits citations.json
  3. Apply:     python3 reinstate_citations.py --deck deck.pptx --edits citations.json --apply
"""
import argparse
import json
import shutil
import sys
import time
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def walk(shapes):
    for s in shapes:
        if s.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk(s.shapes)
        else:
            yield s


def append_cite(run, cite):
    t = run.text
    if t.rstrip().endswith("."):
        i = t.rstrip().rfind(".")
        run.text = t.rstrip()[:i] + " " + cite + "."
    else:
        run.text = t.rstrip() + " " + cite


def safe_backup(deck, phase):
    """Copy deck to <deck>_BACKUP_<phase>.pptx, never clobbering an existing backup."""
    base = deck.rsplit(".pptx", 1)[0] + f"_BACKUP_{phase}"
    path = base + ".pptx"
    import os
    if os.path.exists(path):
        path = base + "_" + time.strftime("%Y%m%d-%H%M%S") + ".pptx"
    shutil.copyfile(deck, path)
    return path


def load_edits(path):
    raw = json.load(open(path, encoding="utf-8"))
    return [(int(s), sub, cite) for s, sub, cite in raw]


def main():
    ap = argparse.ArgumentParser(description="Append dropped citations by unique substring.")
    ap.add_argument("--deck", required=True)
    ap.add_argument("--edits", required=True, help="JSON list of [slide, substring, citation]")
    ap.add_argument("--apply", action="store_true", help="write changes (default: dry run)")
    args = ap.parse_args()

    edits = load_edits(args.edits)
    p = Presentation(args.deck)
    slides = list(p.slides)
    problems = skipped = 0

    for sl, sub, cite in edits:
        if not 1 <= sl <= len(slides):
            print(f"!!! S{sl} | slide out of range (deck has {len(slides)})")
            problems += 1
            continue
        s = slides[sl - 1]
        matches = []
        for sh in walk(s.shapes):
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                if sub in "".join(r.text for r in para.runs):
                    matches.append(para)
        if len(matches) != 1:
            print(f"!!! S{sl} | {sub!r} matched {len(matches)} (need exactly 1)")
            problems += 1
            continue
        para_text = "".join(r.text for r in matches[0].runs)
        if cite in para_text:
            print(f"SKIP S{sl} -> {cite} (already present)")
            skipped += 1
            continue
        if args.apply and matches[0].runs:
            append_cite(matches[0].runs[-1], cite)
        print(f"OK  S{sl} -> {cite}")

    applied = len(edits) - problems - skipped
    if not args.apply:
        print(f"\nDRY RUN: {applied}/{len(edits)} would be applied "
              f"({skipped} already present, {problems} problems). "
              f"Re-run with --apply to write.")
        return
    if problems:
        print("\nNOT SAVED -- fix the ambiguous/zero-match/out-of-range edits first.")
        sys.exit(1)
    if applied == 0:
        print(f"\nNothing to do ({skipped} already present). Deck unchanged.")
        return
    backup = safe_backup(args.deck, "citations")
    p.save(args.deck)
    print(f"\nSAVED {args.deck}  ({applied} applied, {skipped} skipped; backup: {backup})")


if __name__ == "__main__":
    main()

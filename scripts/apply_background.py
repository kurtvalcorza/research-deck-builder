#!/usr/bin/env python3
"""Apply a full-bleed background image to every slide of a .pptx (Repair track).

Sets each slide's <p:bg> to a stretched picture fill referencing --img. Works on decks
built by build_deck_template.js (which set a per-slide solid background) -- the solid fill
is replaced. Content, citations, and speaker notes are untouched.

Safety: dry run by default (reports what it would change); --apply writes, auto-backing
up to <deck>_BACKUP_bg.pptx first (an existing backup is never overwritten).

Usage:
  python3 apply_background.py --deck deck.pptx --img assets/background.jpeg          # dry run
  python3 apply_background.py --deck deck.pptx --img assets/background.jpeg --apply
"""
import argparse
import os
import shutil
import sys
import time
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

BGXML = ('<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
         'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
         'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
         '<p:bgPr><a:blipFill><a:blip r:embed="%s"/><a:stretch><a:fillRect/></a:stretch>'
         '</a:blipFill><a:effectLst/></p:bgPr></p:bg>')


def safe_backup(deck, phase="bg"):
    base = deck.rsplit(".pptx", 1)[0] + f"_BACKUP_{phase}"
    path = base + ".pptx"
    if os.path.exists(path):
        path = base + "_" + time.strftime("%Y%m%d-%H%M%S") + ".pptx"
    shutil.copyfile(deck, path)
    return path


def set_bg(slide, img):
    # get_or_add_image_part returns (image_part, rId)
    _, rId = slide.part.get_or_add_image_part(img)
    cSld = slide._element.find(qn('p:cSld'))
    old = cSld.find(qn('p:bg'))
    if old is not None:
        cSld.remove(old)
    cSld.insert(0, parse_xml(BGXML % rId))   # <p:bg> must be the first child of <p:cSld>


def main():
    ap = argparse.ArgumentParser(description="Set a full-bleed background image on every slide.")
    ap.add_argument("--deck", required=True)
    ap.add_argument("--img", required=True, help="background image (16:9, e.g. assets/background.jpeg)")
    ap.add_argument("--apply", action="store_true", help="write changes (default: dry run)")
    args = ap.parse_args()

    if not os.path.isfile(args.img):
        sys.exit(f"ERROR: image not found: {args.img}")
    p = Presentation(args.deck)
    n = len(p.slides)
    if not args.apply:
        print(f"DRY RUN: would replace the background of {n} slides in {args.deck} "
              f"with {args.img}. Re-run with --apply to write.")
        return
    backup = safe_backup(args.deck)
    for s in p.slides:
        set_bg(s, args.img)
    p.save(args.deck)
    print(f"Background applied to {n} slides in {args.deck} (backup: {backup})")


if __name__ == "__main__":
    main()

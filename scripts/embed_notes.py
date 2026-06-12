#!/usr/bin/env python3
"""Embed a per-slide speaker script (JSON) into a deck's Notes pane, and report word counts.

Script JSON: {"1": "script for slide 1", "2": "...", ...} keyed by slide number (1-based).
Avoid double-quotes inside the JSON string values; use single quotes for quoted phrases.
See references/script_template.json for the shape and the house style.

Subcommands:
  embed   : set each slide's notes from --script (overwrites existing notes; backs up first).
  topup   : insert a short sentence BEFORE the final 'The takeaway:' sentence of given slides
            (from --topups JSON {"9": " extra sentence."}), to lengthen ones under the floor.
  report  : print per-slide word counts and ~seconds, and flag any outside the word band.
  export  : dump the deck's CURRENT notes to a script JSON (--out). Run this after any
            topup or hand edit so the JSON stays the source of truth -- a rebuild bakes
            notes from the JSON, so stale JSON silently reverts in-deck edits.

The word band (default 180-210 ~= 90-120s at ~2.3 words/sec, ~138 wpm) applies to content
slides. The title and closing slides are intentionally short (house style: 30-60 words),
so report exempts the first and last slide by default (--exempt to override).

Safety: embed/topup auto-backup to <deck>_BACKUP_notes.pptx and never overwrite an
existing backup (a timestamped name is used instead).

Usage:
  python3 embed_notes.py embed  --deck deck.pptx --script m01_script.json
  python3 embed_notes.py report --deck deck.pptx
  python3 embed_notes.py topup  --deck deck.pptx --topups topups.json
  python3 embed_notes.py export --deck deck.pptx --out m01_script.json
"""
import argparse
import json
import os
import shutil
import sys
import time
from pptx import Presentation

WPS = 2.3          # words/sec estimate for the ~seconds readout (~138 wpm)
MIN_W, MAX_W = 180, 210


def parse_exempt(spec, n_slides):
    """'first,last' / '1,15' / '' -> set of exempt slide numbers."""
    out = set()
    for tok in (spec or "").split(","):
        tok = tok.strip().lower()
        if not tok:
            continue
        if tok == "first":
            out.add(1)
        elif tok == "last":
            out.add(n_slides)
        else:
            out.add(int(tok))
    return out


def report(deck, exempt_spec="first,last"):
    p = Presentation(deck)
    slides = list(p.slides)
    exempt = parse_exempt(exempt_spec, len(slides))
    for i, s in enumerate(slides, 1):
        t = s.notes_slide.notes_text_frame.text if s.has_notes_slide else ""
        w = len(t.split())
        tag = ""
        if w == 0:
            tag = "  <-- EMPTY"
        elif i in exempt:
            tag = "  (exempt)"
        elif not MIN_W <= w <= MAX_W:
            tag = "  <-- outside band"
        print(f"S{i}: {w} words  (~{round(w / WPS)}s){tag}")


def safe_backup(deck, phase="notes"):
    """Copy deck to <deck>_BACKUP_<phase>.pptx, never clobbering an existing backup."""
    base = deck.rsplit(".pptx", 1)[0] + f"_BACKUP_{phase}"
    path = base + ".pptx"
    if os.path.exists(path):
        path = base + "_" + time.strftime("%Y%m%d-%H%M%S") + ".pptx"
    shutil.copyfile(deck, path)
    return path


def check_keys(keys, n_slides, what):
    bad = [k for k in keys if not 1 <= int(k) <= n_slides]
    if bad:
        sys.exit(f"ERROR: {what} keys out of range for a {n_slides}-slide deck: {bad}")


def embed(args):
    notes = json.load(open(args.script, encoding="utf-8"))
    notes = {k: v for k, v in notes.items() if not k.startswith("_")}   # allow _comment keys
    p = Presentation(args.deck)
    slides = list(p.slides)
    check_keys(notes.keys(), len(slides), "script")
    b = safe_backup(args.deck)
    for k, v in notes.items():
        slides[int(k) - 1].notes_slide.notes_text_frame.text = v
    p.save(args.deck)
    print(f"Embedded {len(notes)} scripts (backup: {b})\n")
    report(args.deck, args.exempt)


def topup(args):
    topups = json.load(open(args.topups, encoding="utf-8"))
    topups = {k: v for k, v in topups.items() if not k.startswith("_")}
    p = Presentation(args.deck)
    slides = list(p.slides)
    check_keys(topups.keys(), len(slides), "topups")
    b = safe_backup(args.deck)
    marker = "The takeaway:"
    for n, extra in topups.items():
        tf = slides[int(n) - 1].notes_slide.notes_text_frame
        t = tf.text
        if marker in t:
            pos = t.rindex(marker)
            t = t[:pos].rstrip() + extra + " " + t[pos:]
        else:
            t = t.rstrip() + extra
        tf.text = t
    p.save(args.deck)
    print(f"Topped up {len(topups)} slides (backup: {b})\n")
    print("Now re-sync the script JSON:  python3 embed_notes.py export "
          f"--deck {args.deck} --out <mNN_script.json>\n")
    report(args.deck, args.exempt)


def export(args):
    p = Presentation(args.deck)
    out = {}
    for i, s in enumerate(p.slides, 1):
        t = s.notes_slide.notes_text_frame.text.strip() if s.has_notes_slide else ""
        if t:
            out[str(i)] = t
    payload = json.dumps(out, ensure_ascii=False, indent=2)
    with open(args.out, "w", encoding="utf-8") as f:
        f.write(payload + "\n")
    print(f"Exported notes for {len(out)} slides -> {args.out}")


def main():
    ap = argparse.ArgumentParser(description="Embed/report/export speaker notes in a .pptx.")
    sub = ap.add_subparsers(dest="cmd", required=True)

    pe = sub.add_parser("embed")
    pe.add_argument("--deck", required=True)
    pe.add_argument("--script", required=True, help="per-slide script JSON")
    pe.add_argument("--exempt", default="first,last", help="band-exempt slides (default first,last)")
    pe.set_defaults(func=embed)

    pt = sub.add_parser("topup")
    pt.add_argument("--deck", required=True)
    pt.add_argument("--topups", required=True, help='JSON {"slide": " extra sentence."}')
    pt.add_argument("--exempt", default="first,last")
    pt.set_defaults(func=topup)

    pr = sub.add_parser("report")
    pr.add_argument("--deck", required=True)
    pr.add_argument("--exempt", default="first,last")
    pr.set_defaults(func=lambda a: report(a.deck, a.exempt))

    px = sub.add_parser("export")
    px.add_argument("--deck", required=True)
    px.add_argument("--out", required=True, help="write current notes to this script JSON")
    px.set_defaults(func=export)

    args = ap.parse_args()
    args.func(args)


if __name__ == "__main__":
    main()

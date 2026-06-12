#!/usr/bin/env python3
"""Programmatic QA gate for a finished/edited deck (Phase 3).

Structural checks (always): slide count, speaker notes present on every slide (and within
the word band), inline citations present, and NO leftover icon-font glyphs (the exact
failure mode the visual rebuild exists to fix).

Outline checks (with --outline <mNN_outline.json>): enforce the Phase 0 content plan --
slide count matches the outline, every outline title appears on its slide, and every
outline citation string appears on its slide. These are HARD failures: the outline is
the source of truth for a (re)build, and this is what makes that claim enforceable.

Fidelity checks (with --source <module.md>): PRESENCE checks against the source text --
  * every surname in every citation on a slide appears in the source,
  * each citation's year appears in the source, and
  * every stat-like number on a slide (85%, 18,000, 147.5 ...) appears in the source --
so a wrong attribution or a transposed digit no longer passes QA silently. NOTE the
limit: presence checking cannot tell a right-surname-wrong-claim attribution; that
pairing stays a human read. Numbered cites ([1]) are resolved to surnames via
--refmap <mNN_refmap.json> (from map_references.py); without a refmap they count for
citation presence but cannot be fidelity-checked.

Exit code is non-zero when a HARD invariant fails, so it can gate a build:
  HARD  : missing/empty notes (when required); slide-count mismatch; icon-font text;
          outline mismatches (count / title / citation) when --outline is given.
  SOFT  : notes outside the word band; content slide with no citation; suspicious
          icon-ligature tokens; fidelity misses. With --strict these also fail the run.
Use --strict for final delivery so fidelity misses fail the build.

The title and closing slides are exempt from the word band and the citation-presence
check by default (house style keeps them short); override with --band-exempt/--cite-exempt.

Usage:
  python3 verify_deck.py --deck deck_REDESIGN.pptx
  python3 verify_deck.py --deck deck_REDESIGN.pptx --baseline deck_BACKUP_notes.pptx
  python3 verify_deck.py --deck deck_REDESIGN.pptx --outline m01_outline.json \\
      --source "01. Module Title.md" --refmap m01_refmap.json --strict
  python3 verify_deck.py --deck deck.pptx --no-require-notes   # citations-only pass
"""
import argparse
import json
import re
import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Icon fonts that render as broken text where they are not installed. Any run using one
# of these is a hard failure: rebuild that slide with vector shapes instead (see SKILL.md).
ICON_FONTS = re.compile(r"material\s*(symbols|icons)|font\s*awesome|glyphicons", re.I)

# Material-Symbols ligatures look like normal snake_case words once the font drops out
# (e.g. 'check_circle', 'arrow_forward'). Flag short all-lowercase underscore tokens that
# sit alone in a run -- a soft signal, since legitimate code samples also use snake_case.
LIGATURE = re.compile(r"^[a-z]{2,}(?:_[a-z]{2,}){1,3}$")

# Inline citation shapes the deck uses: author-date, bare author-name, org-with-year.
#   (Park & Choo, 2024) | (Sun et al., 2025) | (Kulkarni, 2024) | (Nash) | (NIST, 2024)
# The lookahead requires a year somewhere in the parenthetical OR a lowercase letter in
# the first name token, so ALL-CAPS acronyms like (REFINE) / (AI) do NOT count.
CITATION = re.compile(
    r"\((?=[^()]*\d{4}|[A-Z][A-Za-zÀ-ſ.'-]*[a-zà-ſ])"
    r"([A-Z][A-Za-zÀ-ſ.'-]+)"
    r"(?:\s+et al\.|\s+(?:&|and)\s+[A-Z][A-Za-zÀ-ſ.'-]+)?"
    r"(?:\s*,\s*\d{4}[a-z]?)?\)"
)
# Numbered cites like [1] or [2, 5] (IEEE style on slides)
NUMBERED_CITE = re.compile(r"\[(\d+(?:\s*,\s*\d+)*)\]")
# Surname tokens inside a matched citation: must contain a lowercase letter, which
# filters 'et'/'al'/'and' (no leading capital) and years/acronyms.
SURNAME_IN_CITE = re.compile(r"[A-Z][A-Za-zÀ-ſ'-]*[a-zà-ſ][A-Za-zÀ-ſ'-]*")

# Stat-like numbers worth checking against the source: has %, a decimal, a thousands
# comma, or an integer value > 20 (filters out chip/step/index numerals).
NUMBER = re.compile(r"\b\d{1,3}(?:,\d{3})+(?:\.\d+)?%?|\b\d+\.\d+%?|\b\d+%|\b\d+\b")


def walk(shapes):
    for s in shapes:
        if s.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk(s.shapes)
        else:
            yield s


def parse_exempt(spec, n_slides):
    out = set()
    for tok in (spec or "").split(","):
        tok = tok.strip().lower()
        if not tok:
            continue
        if tok == "first":
            out.add(1)
        elif tok == "last":
            out.add(n_slides)
        else:
            out.add(int(tok))
    return out


def stat_numbers(text):
    """Stat-like number tokens in text (citations should be stripped beforehand)."""
    keep = []
    for tok in NUMBER.findall(text):
        plain = tok.rstrip("%").replace(",", "")
        try:
            val = float(plain)
        except ValueError:
            continue
        if "%" in tok or "." in tok or "," in tok or val > 20:
            keep.append(tok)
    return keep


def verify(args):
    p = Presentation(args.deck)
    slides = list(p.slides)
    n = len(slides)
    band_exempt = parse_exempt(args.band_exempt, n)
    cite_exempt = parse_exempt(args.cite_exempt, n)

    source_text = source_flat = None
    if args.source:
        sp = Path(args.source)
        if not sp.is_file():
            sys.exit(f"ERROR: --source not found: {sp}")
        source_text = sp.read_text(encoding="utf-8", errors="replace")
        source_flat = source_text.replace(",", "")     # match 18000 against "18,000" too

    outline = None
    if args.outline:
        op = Path(args.outline)
        if not op.is_file():
            sys.exit(f"ERROR: --outline not found: {op}")
        outline = json.loads(op.read_text(encoding="utf-8"))
        if not outline.get("slides"):
            sys.exit(f"ERROR: --outline {op} has no 'slides' list")

    refmap = None
    if args.refmap:
        rp = Path(args.refmap)
        if not rp.is_file():
            sys.exit(f"ERROR: --refmap not found: {rp}")
        refmap = json.loads(rp.read_text(encoding="utf-8")).get("entries", {})

    hard, soft = [], []

    # --- slide count (precedence: --expect > outline > --baseline) ---
    expected = args.expect
    if expected is None and outline is not None:
        expected = len(outline["slides"])
    if expected is None and args.baseline:
        expected = len(Presentation(args.baseline).slides)
    if expected is not None and n != expected:
        hard.append(f"slide count is {n}, expected {expected}")

    slide_texts = {}                      # kept for the outline checks below
    numbered_without_refmap = False

    for i, s in enumerate(slides, 1):
        # collect text + fonts on this slide
        texts, fonts = [], set()
        for sh in walk(s.shapes):
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                for r in para.runs:
                    if r.text:
                        texts.append(r.text)
                    if r.font.name:
                        fonts.add(r.font.name)
        slide_text = " ".join(texts)
        slide_texts[i] = slide_text

        # --- icon-font (hard) ---
        bad_fonts = [f for f in fonts if ICON_FONTS.search(f or "")]
        if bad_fonts:
            hard.append(f"S{i}: icon font still in use -> {', '.join(sorted(bad_fonts))}")
        for tok in texts:
            t = tok.strip()
            if LIGATURE.match(t) and "Consolas" not in fonts:
                soft.append(f"S{i}: possible icon-ligature artifact text '{t}'")

        # --- notes (hard if required) ---
        notes = s.notes_slide.notes_text_frame.text.strip() if s.has_notes_slide else ""
        wc = len(notes.split())
        if args.require_notes and not notes:
            hard.append(f"S{i}: no speaker notes")
        elif notes and i not in band_exempt and not (args.min_words <= wc <= args.max_words):
            soft.append(f"S{i}: notes {wc} words (band {args.min_words}-{args.max_words})")

        # --- citations present (soft, informational) ---
        cite_strs = [m.group(0) for m in CITATION.finditer(slide_text)]
        numbered_ids = sorted({t.strip() for g in NUMBERED_CITE.findall(slide_text)
                               for t in g.split(",")}, key=int)
        if not cite_strs and not numbered_ids and i not in cite_exempt:
            soft.append(f"S{i}: no inline citation detected")

        # --- fidelity vs source (soft) ---
        if source_text is not None:
            # author-date / bare-name cites: every surname + the year must appear
            for cite in sorted(set(cite_strs)):
                for name in set(SURNAME_IN_CITE.findall(cite)):
                    if name not in source_text:
                        soft.append(f"S{i}: cited surname '{name}' not found in source")
                ym = re.search(r"\d{4}", cite)
                if ym and ym.group(0) not in source_text:
                    soft.append(f"S{i}: cited year '{ym.group(0)}' not found in source")
            # numbered cites: resolve via the refmap, then check those surnames
            if numbered_ids and refmap is not None:
                for nid in numbered_ids:
                    entry = refmap.get(nid)
                    if not entry:
                        soft.append(f"S{i}: cites [{nid}] but refmap has no entry {nid}")
                        continue
                    for name in entry.get("surnames", []):
                        if name not in source_text:
                            soft.append(f"S{i}: [{nid}] surname '{name}' not found in source")
            elif numbered_ids and refmap is None:
                numbered_without_refmap = True
            no_cite_text = NUMBERED_CITE.sub(" ", CITATION.sub(" ", slide_text))
            for num in set(stat_numbers(no_cite_text)):
                plain = num.replace(",", "")
                if num not in source_text and plain not in source_flat:
                    soft.append(f"S{i}: number '{num}' not found in source")

    # --- outline enforcement (HARD): the Phase 0 plan is the source of truth ---
    if outline is not None:
        def norm(t):
            return re.sub(r"\s+", " ", t).strip().casefold()
        for osl in outline["slides"]:
            k = osl.get("n")
            if not isinstance(k, int) or not 1 <= k <= n:
                hard.append(f"outline slide n={k!r} is out of range for a {n}-slide deck")
                continue
            st = norm(slide_texts.get(k, ""))
            title = osl.get("title") or ""
            if title and norm(title) not in st:
                hard.append(f"S{k}: outline title {title!r} not found on slide")
            for c in osl.get("citations") or []:
                if norm(c) not in st:
                    hard.append(f"S{k}: outline citation '{c}' not found on slide")
        # consume the outline's citation_style declaration (see C2/C3)
        if outline.get("citation_style") == "numbered" and refmap is None and source_text is not None:
            soft.append("outline declares numbered citation style — pass --refmap "
                        "mNN_refmap.json so [n] cites can be fidelity-checked")
    if numbered_without_refmap and source_text is not None:
        soft.append("numbered [n] citations found on slides but no --refmap given — "
                    "they count for citation presence but were not fidelity-checked")

    # --- report ---
    print(f"Deck: {args.deck}")
    print(f"Slides: {n}" + (f" (expected {expected})" if expected is not None else ""))
    if args.outline:
        print(f"Outline enforced: {args.outline}")
    if args.source:
        print(f"Source fidelity: checked against {args.source}"
              + (f" (refmap: {args.refmap})" if args.refmap else ""))
    print(f"HARD failures: {len(hard)} | SOFT warnings: {len(soft)}\n")
    for h in hard:
        print(f"  [FAIL] {h}")
    for w in soft:
        print(f"  [warn] {w}")
    if not hard and not soft:
        print("  All checks passed.")

    failed = bool(hard) or (args.strict and bool(soft))
    print("\nRESULT:", "FAIL" if failed else "PASS")
    return 1 if failed else 0


def main():
    ap = argparse.ArgumentParser(description="QA gate for a research deck (.pptx).")
    ap.add_argument("--deck", required=True, help="deck to check")
    ap.add_argument("--source", help="source module .md for fidelity checks (names, years, numbers)")
    ap.add_argument("--outline", help="mNN_outline.json — enforce the Phase 0 plan (count, titles, citations; HARD)")
    ap.add_argument("--refmap", help="mNN_refmap.json — resolve numbered [n] cites for fidelity checks")
    ap.add_argument("--expect", type=int, help="expected slide count (overrides --outline/--baseline)")
    ap.add_argument("--baseline", help="another deck to read the expected slide count from")
    ap.add_argument("--min-words", type=int, default=200, help="notes word floor (default 200)")
    ap.add_argument("--max-words", type=int, default=270, help="notes word ceiling (default 270)")
    ap.add_argument("--band-exempt", default="first,last",
                    help="slides exempt from the word band (default: first,last)")
    ap.add_argument("--cite-exempt", default="first,last",
                    help="slides exempt from the citation-presence check (default: first,last)")
    ap.add_argument("--require-notes", dest="require_notes", action="store_true", default=True,
                    help="fail if any slide lacks notes (default)")
    ap.add_argument("--no-require-notes", dest="require_notes", action="store_false",
                    help="don't require notes (e.g. citations-only pass)")
    ap.add_argument("--strict", action="store_true", help="treat SOFT warnings as failures")
    sys.exit(verify(ap.parse_args()))


if __name__ == "__main__":
    main()
